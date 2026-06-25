#!/usr/bin/env python3
"""
File analyzer — produces a structured digest of a document.
Called by the Node worker as: python3 analyze.py <filepath>
Outputs JSON to stdout.

Supported: PDF (.pdf), DOCX (.docx), PPTX (.pptx)
"""
import json
import sys
import os

def analyze_pdf(filepath):
    import fitz  # PyMuPDF
    doc = fitz.open(filepath)
    pages = []
    total_images = 0
    total_tables = 0
    total_words = 0
    for i, page in enumerate(doc):
        text = page.get_text()
        words = text.split()
        images = page.get_images(full=True)
        # Rough table detection: count lines with multiple tab/column separators
        lines = text.split('\n')
        table_lines = sum(1 for line in lines if line.count('\t') >= 2 or line.count('  ') >= 3)
        has_table = table_lines > 2

        total_images += len(images)
        total_words += len(words)
        if has_table:
            total_tables += 1

        pages.append({
            'page': i + 1,
            'words': len(words),
            'images': len(images),
            'has_table': has_table,
            'heading': lines[0].strip()[:120] if lines and lines[0].strip() else None,
            'text': text,
        })
    doc.close()

    return {
        'type': 'pdf',
        'page_count': len(pages),
        'total_words': total_words,
        'total_images': total_images,
        'pages_with_tables': total_tables,
        'pages': pages,
    }

def analyze_pptx(filepath):
    from pptx import Presentation
    from pptx.util import Inches, Pt
    prs = Presentation(filepath)
    slides = []
    total_words = 0
    total_images = 0
    for i, slide in enumerate(prs.slides):
        texts = []
        images = 0
        title = None
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if line:
                        texts.append(line)
            if hasattr(shape, 'image'):
                images += 1
            if shape.shape_type is not None and hasattr(shape, 'text') and not title:
                # Try to grab title
                try:
                    if shape.placeholder_format and shape.placeholder_format.idx == 0:
                        title = shape.text.strip()[:120]
                except Exception:
                    pass
        if not title and texts:
            title = texts[0][:120]

        words = sum(len(t.split()) for t in texts)
        total_words += words
        total_images += images

        slides.append({
            'slide': i + 1,
            'title': title,
            'words': words,
            'images': images,
            'bullet_count': len(texts),
            'text': '\n'.join(texts),
        })

    return {
        'type': 'pptx',
        'slide_count': len(slides),
        'total_words': total_words,
        'total_images': total_images,
        'slides': slides,
    }

def analyze_docx(filepath):
    from docx import Document
    doc = Document(filepath)
    sections = []
    current_heading = None
    current_text = []
    total_words = 0
    total_images = 0
    total_tables = len(doc.tables)

    for para in doc.paragraphs:
        if para.style and para.style.name.startswith('Heading'):
            # Flush previous section
            if current_heading or current_text:
                text = '\n'.join(current_text)
                words = len(text.split())
                total_words += words
                sections.append({
                    'heading': current_heading or '(no heading)',
                    'words': words,
                    'text': text,
                })
            current_heading = para.text.strip()[:120]
            current_text = []
        else:
            if para.text.strip():
                current_text.append(para.text.strip())

    # Flush last section
    if current_heading or current_text:
        text = '\n'.join(current_text)
        words = len(text.split())
        total_words += words
        sections.append({
            'heading': current_heading or '(no heading)',
            'words': words,
            'text': text,
        })

    # Count images (inline shapes)
    for rel in doc.part.rels.values():
        if "image" in rel.reltype:
            total_images += 1

    return {
        'type': 'docx',
        'section_count': len(sections),
        'total_words': total_words,
        'total_images': total_images,
        'total_tables': total_tables,
        'sections': sections,
    }

def format_digest(filename, analysis):
    """Build a human-readable digest from the structured analysis."""
    lines = [f'=== File: {filename} ===']
    ftype = analysis['type']

    if ftype == 'pdf':
        lines.append(f"Type: PDF | Pages: {analysis['page_count']} | Words: {analysis['total_words']} | Images: {analysis['total_images']} | Pages with tables: {analysis['pages_with_tables']}")
        lines.append('')
        for p in analysis['pages']:
            heading = f" — {p['heading']}" if p['heading'] else ''
            extras = []
            if p['images']:
                extras.append(f"{p['images']} images")
            if p['has_table']:
                extras.append('table detected')
            extra_str = f" [{', '.join(extras)}]" if extras else ''
            lines.append(f"Page {p['page']}{heading}{extra_str}")
            lines.append(p['text'])
            lines.append('')

    elif ftype == 'pptx':
        lines.append(f"Type: PPTX | Slides: {analysis['slide_count']} | Words: {analysis['total_words']} | Images: {analysis['total_images']}")
        lines.append('')
        for s in analysis['slides']:
            title = f" — {s['title']}" if s['title'] else ''
            extras = []
            if s['images']:
                extras.append(f"{s['images']} images")
            extra_str = f" [{', '.join(extras)}]" if extras else ''
            lines.append(f"Slide {s['slide']}{title} ({s['bullet_count']} bullets, {s['words']} words){extra_str}")
            lines.append(s['text'])
            lines.append('')

    elif ftype == 'docx':
        lines.append(f"Type: DOCX | Sections: {analysis['section_count']} | Words: {analysis['total_words']} | Images: {analysis['total_images']} | Tables: {analysis['total_tables']}")
        lines.append('')
        for sec in analysis['sections']:
            lines.append(f"## {sec['heading']} ({sec['words']} words)")
            lines.append(sec['text'])
            lines.append('')

    return '\n'.join(lines)

def main():
    if len(sys.argv) < 2:
        print(json.dumps({'error': 'Usage: analyze.py <filepath>'}))
        sys.exit(1)

    filepath = sys.argv[1]
    filename = os.path.basename(filepath)

    if not os.path.exists(filepath):
        print(json.dumps({'error': f'File not found: {filepath}'}))
        sys.exit(1)

    ext = os.path.splitext(filepath)[1].lower()
    try:
        if ext == '.pdf':
            analysis = analyze_pdf(filepath)
        elif ext == '.pptx':
            analysis = analyze_pptx(filepath)
        elif ext in ('.docx', '.doc'):
            analysis = analyze_docx(filepath)
        else:
            print(json.dumps({'error': f'Unsupported file type: {ext}'}))
            sys.exit(1)

        digest = format_digest(filename, analysis)
        result = {
            'filename': filename,
            'analysis': analysis,
            'digest': digest,
        }
        print(json.dumps(result))
    except Exception as e:
        print(json.dumps({'error': str(e)}))
        sys.exit(1)

if __name__ == '__main__':
    main()

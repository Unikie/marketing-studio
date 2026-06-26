#!/usr/bin/env python3
"""
Pyworker — Tool registry and execution engine.
All I/O is JSON. Tools are Python code stored in SQLite, called by name.
"""
import base64
import json
import os
import sqlite3
import tempfile
import traceback

from flask import Flask, request, jsonify, g
from sqlalchemy import Column, DateTime, Integer, MetaData, Table, Text, create_engine, delete, func, insert, inspect, select, text, update

app = Flask(__name__)

DATABASE_URL = os.environ.get('DATABASE_URL', '')
DB_PATH = os.environ.get('PYWORKER_DB', os.path.join(os.path.dirname(__file__), 'tools.db'))

def get_database_url():
    if DATABASE_URL.startswith('postgres'):
        if DATABASE_URL.startswith('postgres://'):
            return DATABASE_URL.replace('postgres://', 'postgresql+psycopg://', 1)
        if DATABASE_URL.startswith('postgresql://'):
            return DATABASE_URL.replace('postgresql://', 'postgresql+psycopg://', 1)
        return DATABASE_URL
    return f'sqlite:///{DB_PATH}'

DB_DRIVER = 'postgres' if DATABASE_URL.startswith('postgres') else 'sqlite'
engine = create_engine(
    get_database_url(),
    connect_args={'check_same_thread': False} if DB_DRIVER == 'sqlite' else {},
    future=True,
)

metadata = MetaData()
tools_table = Table(
    'tools',
    metadata,
    Column('name', Text, primary_key=True),
    Column('description', Text),
    Column('code', Text, nullable=False),
    Column('params_schema', Text),
    Column('owner', Text, default='system', server_default='system'),
    Column('read_only', Integer, default=0, server_default='0'),
    Column('permissions', Text, default='{}', server_default='{}'),
    Column('created_at', DateTime, server_default=func.now()),
    Column('updated_at', DateTime, server_default=func.now()),
)

# ---- Database ----

def get_db():
    if 'db' not in g:
        g.db = engine.connect()
    return g.db

@app.teardown_appcontext
def close_db(exc):
    db = g.pop('db', None)
    if db:
        db.close()

def row_to_dict(row):
    if not row:
        return None
    result = dict(row)
    for key, value in result.items():
        if hasattr(value, 'isoformat'):
            result[key] = value.isoformat()
    return result

def init_db():
    metadata.create_all(engine)
    if DB_DRIVER == 'sqlite':
        with engine.connect() as db:
            db.execute(text('PRAGMA journal_mode=WAL'))

    columns = [column['name'] for column in inspect(engine).get_columns('tools')]
    if 'params_schema' not in columns:
        with engine.begin() as db:
            db.execute(text('ALTER TABLE tools ADD COLUMN params_schema TEXT'))

# ---- Permission helpers ----

def get_caller():
    """Extract caller identity from request header."""
    return request.headers.get('X-Caller', 'anonymous')

def get_tool_or_404(name):
    row = get_db().execute(select(tools_table).where(tools_table.c.name == name)).mappings().fetchone()
    return row_to_dict(row)

def check_permission(tool, action):
    """Check if caller has permission for action (create/read/update/delete/execute)."""
    caller = get_caller()
    if caller == 'system' or caller == tool.get('owner'):
        return True
    if tool.get('read_only') and action in ('update', 'delete'):
        return False
    perms = json.loads(tool.get('permissions') or '{}')
    user_perms = perms.get(caller, perms.get('*', []))
    return action in user_perms

# ---- Tool CRUD ----

@app.route('/tools', methods=['GET'])
def list_tools():
    rows = get_db().execute(
        select(
            tools_table.c.name,
            tools_table.c.description,
            tools_table.c.params_schema,
            tools_table.c.owner,
            tools_table.c.read_only,
            tools_table.c.created_at,
            tools_table.c.updated_at,
        )
    ).mappings().fetchall()
    return jsonify([row_to_dict(r) for r in rows])

@app.route('/tools', methods=['POST'])
def create_tool():
    data = request.get_json()
    if not data or not data.get('name') or not data.get('code'):
        return jsonify({'error': 'name and code required'}), 400

    name = data['name']
    existing = get_db().execute(select(tools_table.c.name).where(tools_table.c.name == name)).fetchone()
    if existing:
        return jsonify({'error': f'Tool "{name}" already exists'}), 409

    caller = get_caller()
    get_db().execute(
        insert(tools_table).values(
            name=name,
            description=data.get('description', ''),
            code=data['code'],
            params_schema=data.get('params_schema'),
            owner=caller,
            read_only=int(data.get('read_only', False)),
            permissions=json.dumps(data.get('permissions', {})),
        )
    )
    get_db().commit()
    return jsonify({'ok': True, 'name': name}), 201

@app.route('/tools/<name>', methods=['GET'])
def read_tool(name):
    tool = get_tool_or_404(name)
    if not tool:
        return jsonify({'error': 'Not found'}), 404
    if not check_permission(tool, 'read'):
        return jsonify({'error': 'Forbidden'}), 403
    return jsonify(tool)

@app.route('/tools/<name>', methods=['PUT'])
def update_tool(name):
    tool = get_tool_or_404(name)
    if not tool:
        return jsonify({'error': 'Not found'}), 404
    if not check_permission(tool, 'update'):
        return jsonify({'error': 'Forbidden: tool is read-only or insufficient permissions'}), 403

    data = request.get_json()
    if not data:
        return jsonify({'error': 'No data'}), 400

    values = {}
    if 'code' in data:
        values['code'] = data['code']
    if 'description' in data:
        values['description'] = data['description']
    if 'params_schema' in data:
        values['params_schema'] = data['params_schema']
    if values:
        values['updated_at'] = func.now()
        get_db().execute(update(tools_table).where(tools_table.c.name == name).values(**values))
        get_db().commit()
    return jsonify({'ok': True})

@app.route('/tools/<name>', methods=['DELETE'])
def delete_tool(name):
    tool = get_tool_or_404(name)
    if not tool:
        return jsonify({'error': 'Not found'}), 404
    if not check_permission(tool, 'delete'):
        return jsonify({'error': 'Forbidden'}), 403

    get_db().execute(delete(tools_table).where(tools_table.c.name == name))
    get_db().commit()
    return jsonify({'ok': True})

# ---- System: set attributes (owner, read_only, permissions) ----

@app.route('/tools/<name>/attributes', methods=['PUT'])
def set_attributes(name):
    caller = get_caller()
    tool = get_tool_or_404(name)
    if not tool:
        return jsonify({'error': 'Not found'}), 404
    if caller != 'system' and tool.get('owner') != caller:
        return jsonify({'error': 'Only system or owner can set attributes'}), 403

    data = request.get_json()
    if not data:
        return jsonify({'error': 'No data'}), 400

    values = {}
    if 'owner' in data:
        values['owner'] = data['owner']
    if 'read_only' in data:
        values['read_only'] = int(data['read_only'])
    if 'permissions' in data:
        values['permissions'] = json.dumps(data['permissions'])
    if values:
        values['updated_at'] = func.now()
        get_db().execute(update(tools_table).where(tools_table.c.name == name).values(**values))
        get_db().commit()
    return jsonify({'ok': True})

# ---- Tool execution ----

@app.route('/tools/<name>/execute', methods=['POST'])
def execute_tool(name):
    tool = get_tool_or_404(name)
    if not tool:
        return jsonify({'error': f'Tool "{name}" not found'}), 404
    if not check_permission(tool, 'execute'):
        return jsonify({'error': 'Forbidden'}), 403

    params = request.get_json()
    if params is None:
        params = {}

    # Execute tool code in a sandboxed namespace
    namespace = {
        '__builtins__': __builtins__,
        'json': json,
        'base64': base64,
        'tempfile': tempfile,
        'os': os,
    }

    try:
        exec(compile(tool['code'], f'<tool:{name}>', 'exec'), namespace)

        if 'run' not in namespace:
            return jsonify({'error': f'Tool "{name}" does not define a run(params) function'}), 500

        result = namespace['run'](params)

        if not isinstance(result, dict):
            result = {'result': result}

        return jsonify(result)

    except Exception as e:
        return jsonify({'error': str(e), 'traceback': traceback.format_exc()}), 500

# ---- Health ----

@app.route('/health', methods=['GET'])
def health():
    count = get_db().execute(select(func.count()).select_from(tools_table)).scalar_one()
    return jsonify({'status': 'ok', 'tools': count})

# ---- DB export/import (SQL text dump) ----

@app.route('/db/export', methods=['GET'])
def db_export():
    if DB_DRIVER != 'sqlite':
        return '', 501

    db = sqlite3.connect(DB_PATH)
    db.row_factory = sqlite3.Row
    lines = []
    tables = db.execute("SELECT name, sql FROM sqlite_master WHERE type='table' AND name NOT LIKE 'sqlite_%'").fetchall()
    for t in tables:
        lines.append(t['sql'] + ';')
        rows = db.execute(f'SELECT * FROM "{t["name"]}"').fetchall()
        for row in rows:
            cols = row.keys()
            vals = []
            for c in cols:
                v = row[c]
                if v is None:
                    vals.append('NULL')
                elif isinstance(v, (int, float)):
                    vals.append(str(v))
                else:
                    vals.append("'" + str(v).replace("'", "''") + "'")
            col_str = ','.join(f'"{c}"' for c in cols)
            val_str = ','.join(vals)
            lines.append(f'INSERT INTO "{t["name"]}" ({col_str}) VALUES ({val_str});')
    dump = '\n'.join(lines) + '\n'
    db.close()
    return app.response_class(dump, mimetype='text/plain',
                              headers={'Content-Disposition': 'attachment; filename=tools.sql'})

@app.route('/db/import', methods=['POST'])
def db_import():
    if DB_DRIVER != 'sqlite':
        return '', 501

    sql = request.get_data(as_text=True)
    if not sql or len(sql) < 10:
        return jsonify({'error': 'Empty SQL dump'}), 400
    # Close current connection, recreate DB
    db = g.pop('db', None)
    if db:
        db.close()
    engine.dispose()
    if os.path.exists(DB_PATH):
        os.unlink(DB_PATH)
    new_db = sqlite3.connect(DB_PATH)
    new_db.executescript(sql)
    new_db.close()
    return jsonify({'ok': True, 'size': len(sql)})

# ---- Built-in tool: file_analysis ----

TOOL_FILE_ANALYSIS = r'''
import base64
import os
import tempfile

def run(params):
    file_b64 = params.get('file')
    filename = params.get('filename', 'unknown')

    if not file_b64:
        return {'error': 'No file provided (expected base64 in "file" field)'}

    ext = os.path.splitext(filename)[1].lower()
    if ext not in ('.pdf', '.pptx', '.docx', '.doc'):
        return {'error': f'Unsupported file type: {ext}'}

    file_bytes = base64.b64decode(file_b64)

    fd, tmppath = tempfile.mkstemp(suffix=ext)
    try:
        os.write(fd, file_bytes)
        os.close(fd)

        if ext == '.pdf':
            return analyze_pdf(tmppath, filename)
        elif ext == '.pptx':
            return analyze_pptx(tmppath, filename)
        elif ext in ('.docx', '.doc'):
            return analyze_docx(tmppath, filename)
    finally:
        if os.path.exists(tmppath):
            os.unlink(tmppath)


def analyze_pdf(filepath, filename):
    import fitz
    doc = fitz.open(filepath)
    pages = []
    total_images = 0
    total_tables = 0
    total_words = 0
    for i, page in enumerate(doc):
        text = page.get_text()
        words = text.split()
        images = page.get_images(full=True)
        lines = text.split('\n')
        table_lines = sum(1 for line in lines if line.count('\t') >= 2 or line.count('  ') >= 3)
        has_table = table_lines > 2
        total_images += len(images)
        total_words += len(words)
        if has_table:
            total_tables += 1
        pages.append({
            'page': i + 1,
            'words': len(words),
            'images': len(images),
            'has_table': has_table,
            'heading': lines[0].strip()[:120] if lines and lines[0].strip() else None,
            'text': [line for line in lines if line.strip()],
        })
    doc.close()
    return {
        'filename': filename,
        'type': 'pdf',
        'page_count': len(pages),
        'total_words': total_words,
        'total_images': total_images,
        'pages_with_tables': total_tables,
        'pages': pages,
    }


def analyze_pptx(filepath, filename):
    from pptx import Presentation
    prs = Presentation(filepath)
    slides = []
    total_words = 0
    total_images = 0
    for i, slide in enumerate(prs.slides):
        texts = []
        images = 0
        title = None
        for shape in slide.shapes:
            if shape.has_text_frame:
                for para in shape.text_frame.paragraphs:
                    line = para.text.strip()
                    if line:
                        texts.append(line)
            if hasattr(shape, 'image'):
                images += 1
            if not title:
                try:
                    if shape.placeholder_format and shape.placeholder_format.idx == 0:
                        title = shape.text.strip()[:120]
                except Exception:
                    pass
        if not title and texts:
            title = texts[0][:120]
        words = sum(len(t.split()) for t in texts)
        total_words += words
        total_images += images
        slides.append({
            'slide': i + 1,
            'title': title,
            'words': words,
            'images': images,
            'bullet_count': len(texts),
            'text': texts,
        })
    return {
        'filename': filename,
        'type': 'pptx',
        'slide_count': len(slides),
        'total_words': total_words,
        'total_images': total_images,
        'slides': slides,
    }


def analyze_docx(filepath, filename):
    from docx import Document
    doc = Document(filepath)
    sections = []
    current_heading = None
    current_text = []
    total_words = 0
    total_images = 0
    total_tables = len(doc.tables)

    for para in doc.paragraphs:
        if para.style and para.style.name.startswith('Heading'):
            if current_heading or current_text:
                words = sum(len(t.split()) for t in current_text)
                total_words += words
                sections.append({
                    'heading': current_heading or '(no heading)',
                    'words': words,
                    'text': current_text,
                })
            current_heading = para.text.strip()[:120]
            current_text = []
        else:
            if para.text.strip():
                current_text.append(para.text.strip())

    if current_heading or current_text:
        words = sum(len(t.split()) for t in current_text)
        total_words += words
        sections.append({
            'heading': current_heading or '(no heading)',
            'words': words,
            'text': current_text,
        })

    for rel in doc.part.rels.values():
        if "image" in rel.reltype:
            total_images += 1

    return {
        'filename': filename,
        'type': 'docx',
        'section_count': len(sections),
        'total_words': total_words,
        'total_images': total_images,
        'total_tables': total_tables,
        'sections': sections,
    }
'''

# ---- Startup: seed built-in tools ----

TOOL_FILE_ANALYSIS_SCHEMA = json.dumps({
    "type": "object",
    "properties": {
        "file": {"type": "string", "description": "Base64-encoded file content"},
        "filename": {"type": "string", "description": "Original filename with extension (e.g. report.pdf)"}
    },
    "required": ["file", "filename"]
})

def seed_builtin_tools():
    # Seed: lowercase tool
    TOOL_LOWERCASE_CODE = '''def run(params):
    text = params.get('text', '')
    if not text:
        return {'error': 'No text provided (expected "text" field)'}
    return {'result': text.lower()}
'''
    TOOL_LOWERCASE_SCHEMA = json.dumps({
        "type": "object",
        "properties": {
            "text": {"type": "string", "description": "The text to convert to lowercase"}
        },
        "required": ["text"]
    })

    with engine.begin() as db:
        existing = db.execute(select(tools_table.c.name).where(tools_table.c.name == 'file_analysis')).fetchone()
        if not existing:
            db.execute(
                insert(tools_table).values(
                    name='file_analysis',
                    description='Analyze PDF, PPTX, DOCX files. Input: {file: base64, filename: str}',
                    code=TOOL_FILE_ANALYSIS,
                    params_schema=TOOL_FILE_ANALYSIS_SCHEMA,
                    owner='system',
                    read_only=1,
                    permissions=json.dumps({'*': ['read', 'execute']}),
                )
            )
        else:
            db.execute(
                update(tools_table)
                .where(tools_table.c.name == 'file_analysis')
                .values(params_schema=TOOL_FILE_ANALYSIS_SCHEMA, updated_at=func.now())
            )

        existing_lc = db.execute(select(tools_table.c.name).where(tools_table.c.name == 'lowercase')).fetchone()
        if not existing_lc:
            db.execute(
                insert(tools_table).values(
                    name='lowercase',
                    description='Convert all text to lowercase',
                    code=TOOL_LOWERCASE_CODE,
                    params_schema=TOOL_LOWERCASE_SCHEMA,
                    owner='system',
                    read_only=0,
                    permissions=json.dumps({'*': ['read', 'execute']}),
                )
            )

# ---- Main ----

if __name__ == '__main__':
    init_db()
    seed_builtin_tools()
    port = int(os.environ.get('PYWORKER_PORT', '3002'))
    print(f'[pyworker] Tool engine on http://localhost:{port}')
    print(f'[pyworker] DB: {DB_DRIVER} ({DB_PATH if DB_DRIVER == "sqlite" else "DATABASE_URL"})')
    app.run(host='0.0.0.0', port=port, debug=False)
